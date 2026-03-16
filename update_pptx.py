"""Update caddi_drawer_proposal.pptx: add architecture slides after slide 11."""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

prs = Presentation("caddi_drawer_proposal.pptx")
layout = prs.slide_layouts[0]  # DEFAULT

# Style constants (matched from existing slides)
TEAL = RGBColor(0x0D, 0x73, 0x77)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0x66, 0x66, 0x66)
MUTED = RGBColor(0x99, 0x99, 0x99)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_BG = RGBColor(0xFD, 0xE8, 0xE8)
GREEN_BG = RGBColor(0xE8, 0xF5, 0xE9)
BLUE_BG = RGBColor(0xE3, 0xF2, 0xFD)
TEAL_BG = RGBColor(0xE0, 0xF2, 0xF1)
ORANGE_BG = RGBColor(0xFF, 0xF3, 0xE0)
RED = RGBColor(0xC6, 0x28, 0x28)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
BLUE = RGBColor(0x15, 0x65, 0xC0)

LEFT_MARGIN = Emu(731520)
CONTENT_W = Emu(7680960)
FONT_NAME = "Calibri"


def add_textbox(slide, left, top, width, height, text, font_size=12,
                bold=False, color=DARK, alignment=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    if alignment:
        p.alignment = alignment
    return txBox


def add_line(slide, left, top, width):
    """Add a thin horizontal line."""
    shape = slide.shapes.add_shape(
        1, left, top, width, Emu(0)  # 1 = rectangle, 0 height = line
    )
    shape.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    shape.line.width = Pt(1)
    shape.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color, text="",
                     font_size=10, font_color=DARK, bold=False, border_color=None):
    shape = slide.shapes.add_shape(
        5, left, top, width, height  # 5 = rounded rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = FONT_NAME
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = font_color
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow_down(slide, cx, top, length):
    """Add a downward arrow using a simple line with arrow."""
    shape = slide.shapes.add_connector(
        1, cx, top, cx, top + length  # straight connector
    )
    shape.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
    shape.line.width = Pt(1.5)
    return shape


# =========================================================================
# SLIDE 12: Architecture Overview
# =========================================================================

slide_arch = prs.slides.add_slide(layout)

# Title
add_textbox(slide_arch, LEFT_MARGIN, Emu(365760), Emu(8229600), Emu(457200),
            "How the Engine Works", font_size=24, bold=True, color=DARK)

# Subtitle line
add_line(slide_arch, LEFT_MARGIN, Emu(868680), CONTENT_W)

add_textbox(slide_arch, LEFT_MARGIN, Emu(950000), Emu(7680960), Emu(280000),
            "From natural language input to data-backed supplier recommendation in 4 steps",
            font_size=15, bold=True, color=DARK)

# Step boxes - vertical flow
y_start = Emu(1380000)
box_h = Emu(520000)
box_w = Emu(7200000)
gap = Emu(580000)
num_w = Emu(340000)
text_left = Emu(1200000)
text_w = Emu(7200000)

steps = [
    ("1", "User types natural language",
     '"200 aluminum heat exchangers for the data center project"'),
    ("2", "FAISS vector search matches to part category",
     "Sentence-transformers (all-MiniLM-L6-v2) embeds the query. FAISS finds the nearest part category (HX). Scales to hundreds of categories without keyword maps."),
    ("3", "TCO scoring ranks all suppliers for that category",
     "effective_cost = unit_price x (1 + rejection_rate x 0.5 + late_rate x 0.2). Lowest true cost wins. Specialist tiebreaker from team notes."),
    ("4", "Returns 3 results: Recommended, Alternative, Avoid",
     "Each card shows unit price, true cost (TCO), on-time %, rejection %, and team notes. User can ask follow-up questions via Claude chat."),
]

for i, (num, title, desc) in enumerate(steps):
    y = y_start + i * gap

    # Number circle
    add_rounded_rect(slide_arch, LEFT_MARGIN, y, num_w, Emu(280000),
                     TEAL, num, font_size=13, font_color=WHITE, bold=True)

    # Title
    add_textbox(slide_arch, text_left, y - Emu(10000), text_w, Emu(250000),
                title, font_size=13, bold=True, color=DARK)

    # Description
    add_textbox(slide_arch, text_left, y + Emu(220000), text_w, Emu(250000),
                desc, font_size=10, color=LIGHT_GRAY)


# =========================================================================
# SLIDE 13: TCO Formula
# =========================================================================

slide_tco = prs.slides.add_slide(layout)

add_textbox(slide_tco, LEFT_MARGIN, Emu(365760), Emu(8229600), Emu(457200),
            "TCO: Why the Cheapest Quote Isn't the Cheapest Supplier",
            font_size=24, bold=True, color=DARK)

add_line(slide_tco, LEFT_MARGIN, Emu(868680), CONTENT_W)

# Formula box
add_rounded_rect(slide_tco, LEFT_MARGIN, Emu(980000), Emu(8000000), Emu(400000),
                 ORANGE_BG,
                 "effective_cost  =  unit_price  x  ( 1  +  rejection_rate x 0.5  +  late_rate x 0.2 )",
                 font_size=14, font_color=DARK, bold=True,
                 border_color=RGBColor(0xFF, 0x98, 0x00))

# Factor explanations
add_textbox(slide_tco, LEFT_MARGIN, Emu(1520000), Emu(4000000), Emu(200000),
            "Rejection rate x 0.5", font_size=12, bold=True, color=RED)
add_textbox(slide_tco, LEFT_MARGIN, Emu(1720000), Emu(4000000), Emu(250000),
            "Each rejected part costs ~50% of unit price to rework or reorder",
            font_size=10, color=LIGHT_GRAY)

add_textbox(slide_tco, Emu(4800000), Emu(1520000), Emu(4000000), Emu(200000),
            "Late rate x 0.2", font_size=12, bold=True, color=RED)
add_textbox(slide_tco, Emu(4800000), Emu(1720000), Emu(4000000), Emu(250000),
            "Late deliveries add ~20% in expediting, idle lines, missed deadlines",
            font_size=10, color=LIGHT_GRAY)

# Comparison section
add_line(slide_tco, LEFT_MARGIN, Emu(2150000), CONTENT_W)

add_textbox(slide_tco, LEFT_MARGIN, Emu(2220000), Emu(8000000), Emu(250000),
            "Example: Heat Exchangers (HX)", font_size=14, bold=True, color=DARK)

# QuickFab box (red)
qf_left = LEFT_MARGIN
qf_top = Emu(2550000)
add_rounded_rect(slide_tco, qf_left, qf_top, Emu(3700000), Emu(1400000),
                 RED_BG, "", border_color=RGBColor(0xEF, 0x53, 0x50))

add_textbox(slide_tco, Emu(900000), Emu(2600000), Emu(3400000), Emu(200000),
            "QuickFab Industries", font_size=12, bold=True, color=RED)
add_textbox(slide_tco, Emu(900000), Emu(2830000), Emu(3400000), Emu(180000),
            "Unit price: $1,473", font_size=10, color=GRAY)
add_textbox(slide_tco, Emu(900000), Emu(3020000), Emu(3400000), Emu(180000),
            "Rejection: 16.3%  |  Late: 62.5%", font_size=10, color=GRAY)
add_textbox(slide_tco, Emu(900000), Emu(3250000), Emu(3400000), Emu(250000),
            "True cost: $1,777/unit", font_size=13, bold=True, color=RED)
add_textbox(slide_tco, Emu(900000), Emu(3500000), Emu(3400000), Emu(200000),
            "Looks cheapest. Actually most expensive.", font_size=9, color=MUTED)

# Apex box (green)
ax_left = Emu(4800000)
add_rounded_rect(slide_tco, ax_left, qf_top, Emu(3700000), Emu(1400000),
                 GREEN_BG, "", border_color=RGBColor(0x66, 0xBB, 0x6A))

add_textbox(slide_tco, Emu(4970000), Emu(2600000), Emu(3400000), Emu(200000),
            "Apex Manufacturing", font_size=12, bold=True, color=GREEN)
add_textbox(slide_tco, Emu(4970000), Emu(2830000), Emu(3400000), Emu(180000),
            "Unit price: $1,627", font_size=10, color=GRAY)
add_textbox(slide_tco, Emu(4970000), Emu(3020000), Emu(3400000), Emu(180000),
            "Rejection: 2.6%  |  Late: 17.5%", font_size=10, color=GRAY)
add_textbox(slide_tco, Emu(4970000), Emu(3250000), Emu(3400000), Emu(250000),
            "True cost: $1,648/unit", font_size=13, bold=True, color=GREEN)
add_textbox(slide_tco, Emu(4970000), Emu(3500000), Emu(3400000), Emu(200000),
            "Higher sticker price. Lower true cost.", font_size=9, color=MUTED)

# Bottom callout
add_rounded_rect(slide_tco, Emu(2000000), Emu(4200000), Emu(5500000), Emu(350000),
                 TEAL_BG,
                 "Apex saves $129/unit vs QuickFab despite $154 higher sticker price",
                 font_size=11, font_color=TEAL, bold=True,
                 border_color=TEAL)


# =========================================================================
# SLIDE 14: Tech Stack
# =========================================================================

slide_tech = prs.slides.add_slide(layout)

add_textbox(slide_tech, LEFT_MARGIN, Emu(365760), Emu(8229600), Emu(457200),
            "Tech Stack & Architecture", font_size=24, bold=True, color=DARK)

add_line(slide_tech, LEFT_MARGIN, Emu(868680), CONTENT_W)

# Architecture boxes - 3 columns
col_w = Emu(2500000)
col_gap = Emu(200000)
col_h = Emu(2000000)
col_y = Emu(1050000)

# Column 1: Data Layer
c1_left = LEFT_MARGIN
add_rounded_rect(slide_tech, c1_left, col_y, col_w, col_h,
                 TEAL_BG, "", border_color=TEAL)

add_textbox(slide_tech, Emu(900000), Emu(1100000), Emu(2200000), Emu(220000),
            "Data Layer", font_size=13, bold=True, color=TEAL,
            alignment=PP_ALIGN.CENTER)

items_data = [
    "500 purchase orders",
    "200 quality inspections",
    "92 RFQ quotes",
    "Unstructured team notes",
    "",
    "Entity resolution",
    "Inspection-to-supplier join",
    "TCO scoring per supplier",
]
y_item = Emu(1350000)
for item in items_data:
    if item == "":
        y_item += Emu(80000)
        continue
    add_textbox(slide_tech, Emu(900000), y_item, Emu(2200000), Emu(170000),
                item, font_size=9, color=GRAY)
    y_item += Emu(160000)

# Column 2: Backend
c2_left = Emu(c1_left + col_w + col_gap)
add_rounded_rect(slide_tech, c2_left, col_y, col_w, col_h,
                 BLUE_BG, "", border_color=BLUE)

add_textbox(slide_tech, Emu(3550000), Emu(1100000), Emu(2200000), Emu(220000),
            "Backend", font_size=13, bold=True, color=BLUE,
            alignment=PP_ALIGN.CENTER)

items_backend = [
    "Python + Flask",
    "pandas (in-memory data)",
    "FAISS vector index",
    "sentence-transformers",
    "(all-MiniLM-L6-v2)",
    "",
    "4 API endpoints",
    "No database needed",
]
y_item = Emu(1350000)
for item in items_backend:
    if item == "":
        y_item += Emu(80000)
        continue
    add_textbox(slide_tech, Emu(3550000), y_item, Emu(2200000), Emu(170000),
                item, font_size=9, color=GRAY)
    y_item += Emu(160000)

# Column 3: AI / Chat
c3_left = Emu(c2_left + col_w + col_gap)
add_rounded_rect(slide_tech, c3_left, col_y, col_w, col_h,
                 RGBColor(0xED, 0xE7, 0xF6), "", border_color=RGBColor(0x7E, 0x57, 0xC2))

purple = RGBColor(0x7E, 0x57, 0xC2)
add_textbox(slide_tech, Emu(6200000), Emu(1100000), Emu(2200000), Emu(220000),
            "AI Layer", font_size=13, bold=True, color=purple,
            alignment=PP_ALIGN.CENTER)

items_ai = [
    "Anthropic Claude API",
    "claude-sonnet-4-20250514",
    "",
    "System prompt includes:",
    "  All supplier scores",
    "  Team notes",
    "  Current recommendation",
    "",
    "Follow-up Q&A with data",
]
y_item = Emu(1350000)
for item in items_ai:
    if item == "":
        y_item += Emu(80000)
        continue
    add_textbox(slide_tech, Emu(6200000), y_item, Emu(2200000), Emu(170000),
                item, font_size=9, color=GRAY)
    y_item += Emu(160000)

# Bottom: key design decisions
add_line(slide_tech, LEFT_MARGIN, Emu(3250000), CONTENT_W)

add_textbox(slide_tech, LEFT_MARGIN, Emu(3350000), Emu(8000000), Emu(200000),
            "Key Design Decisions", font_size=12, bold=True, color=MUTED)

decisions = [
    "FAISS for part matching    Scales to hundreds of categories. No keyword map maintenance.",
    "TCO as single ranking metric    One number anyone can understand. No arbitrary weight tuning.",
    "Claude for follow-ups only    Recommendation is deterministic from data. AI handles the conversation, not the decision.",
]
y_d = Emu(3600000)
for d in decisions:
    add_textbox(slide_tech, LEFT_MARGIN, y_d, Emu(8200000), Emu(200000),
                d, font_size=9, color=LIGHT_GRAY)
    y_d += Emu(220000)


# =========================================================================
# Now reorder: move new slides (12, 13, 14) before the old slide 12 ("Thank you")
# Slides are currently: [0..11(old thank you), 12(arch), 13(tco), 14(tech)]
# We want: [0..10(prototype), 12(arch), 13(tco), 14(tech), 11(thank you)]
# =========================================================================

slide_list = prs.slides._sldIdLst
slides = list(slide_list)
# Current order: indices 0-14
# slides[11] = old "Thank you" (was slide 12)
# slides[12] = new Architecture
# slides[13] = new TCO
# slides[14] = new Tech Stack
# We want: move slides[11] to the end
thank_you = slides[11]
slide_list.remove(thank_you)
slide_list.append(thank_you)

prs.save("caddi_drawer_proposal.pptx")
print("Done! Updated caddi_drawer_proposal.pptx")
print("New slide order:")
for i, slide in enumerate(prs.slides):
    title_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.paragraphs[0].text.strip()
            if t:
                title_text = t
                break
    print(f"  Slide {i+1}: {title_text[:60]}")
